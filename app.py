import os
import time
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import csv
import numpy as np
import customtkinter
import tkinter as tk
from tkinter import END, messagebox
from tkinter.filedialog import askopenfilenames
from docx import Document
from PyPDF2 import PdfReader
from fpdf import FPDF
from difflib import SequenceMatcher
from PIL import Image, ImageDraw, ImageFont, ImageTk
from sklearn.cluster import MiniBatchKMeans
from sklearn.metrics import silhouette_score
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer
from joblib import Parallel, delayed
from collections import defaultdict
from bs4 import BeautifulSoup  # For HTML
from striprtf.striprtf import rtf_to_text
from odf import text, teletype  # For ODT
from odf.opendocument import load as odf_load  # For ODT
from ebooklib import epub  # For EPUB
import mobi  # For MOBI (requires kindleunpack)
import textract  # For extracting text from various formats, including LaTeX
from openpyxl import load_workbook
import ezodf
from pptx import Presentation
import subprocess
import xml.etree.ElementTree as ET
import re
import sys

# Gets the correct file path whether running as .py or .exe
def get_resource_path(relative_path):
    if getattr(sys, 'frozen', False):  # Running as .exe
        base_path = sys._MEIPASS
    else:  # Running as .py script
        base_path = os.path.dirname(__file__)
    
    return os.path.join(base_path, relative_path)

# Correct file paths for icons and folder
sun_icon = get_resource_path("sun.png")
moon_icon = get_resource_path("moon.png")
nltk_data_path = get_resource_path("nltk_data")

# Define the path to the local NLTK data directory within the project
nltk_data_dir = os.path.join(os.getcwd(), "nltk_data")

# Add the local directory to NLTK's search path
nltk.data.path.append(nltk_data_dir)

# Function to check if a package is already available
def is_package_available(package_name):
    try:
        if package_name == "punkt":
            nltk.data.find("tokenizers/punkt")
        elif package_name == "stopwords":
            nltk.data.find("corpora/stopwords")
        return True
    except LookupError:
        return False

# Helper functions to manage progress bars dynamically
def show_progress_bar(parent, row, column, width=400, height=10):
    """Show a progress bar at the specified grid position."""
    progress = customtkinter.CTkProgressBar(parent, width=width, height=height, mode="determinate", progress_color="#4CAF50")
    progress.grid(row=row, column=column, padx=(100, 20), pady=(0, 10), sticky="ew")  # Equal padding on both sides
    progress.set(0)
    return progress

def hide_progress_bar(progress):
    """Hide the progress bar by removing it from the GUI."""
    if progress:
        progress.grid_remove()

# Ensure required packages are available
required_packages = ["punkt", "stopwords"]
for package in required_packages:
    if not is_package_available(package):
        print(f"NLTK package '{package}' is missing from the local directory.")
    else:
        print(f"NLTK package '{package}' is already available.")

# Initialize Sentence Transformer model
MODEL = SentenceTransformer('all-MiniLM-L6-v2')

# Create necessary folders
for folder in ["Pending", "Reports"]:
    os.makedirs(folder, exist_ok=True)

# File conversion functions
def convert_markdown_to_txt(md_file, txt_file):
    """Convert Markdown (.md or .markdown) to plain text."""
    try:
        with open(md_file, "r", encoding="utf-8") as md:
            content = md.read()
            text_content = re.sub(r'\*|\_|\#|!\[.*?\]\(.*?\)|\[(.*?)\]\(.*?\)', '', content)
            text_content = re.sub(r'\n{2,}', '\n', text_content).strip()
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content)
    except Exception as e:
        print(f"Error converting Markdown file: {e}")

def convert_xml_to_txt(xml_file, txt_file):
    """Convert XML to plain text."""
    try:
        tree = ET.parse(xml_file)
        root = tree.getroot()

        def extract_text(element):
            text = element.text.strip() if element.text else ""
            for child in element:
                text += " " + extract_text(child)
            return text.strip()

        text_content = extract_text(root)
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content)
    except Exception as e:
        print(f"Error converting XML file: {e}")

def convert_odp_to_txt(odp_file, txt_file):
    """Convert ODP to plain text using unoconv."""
    try:
        subprocess.run(["unoconv", "-f", "txt", "-o", txt_file, odp_file], check=True)
    except Exception as e:
        print(f"Error converting ODP file: {e}")

def convert_pptx_to_txt(pptx_file, txt_file):
    """Convert PPTX to plain text."""
    try:
        presentation = Presentation(pptx_file)
        text_content = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content.strip())
    except Exception as e:
        print(f"Error converting PPTX file: {e}")

def convert_ods_to_txt(ods_file, txt_file):
    """Convert ODS to plain text."""
    try:
        doc = ezodf.opendoc(ods_file)
        text_content = ""
        for sheet in doc.sheets:
            text_content += f"Sheet: {sheet.name}\n"
            for row in sheet.rows():
                text_content += "\t".join([str(cell.value) if cell.value is not None else "" for cell in row]) + "\n"
            text_content += "\n"
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content)
    except Exception as e:
        print(f"Error converting ODS file: {e}")

def convert_csv_to_txt(csv_file, txt_file):
    """Convert CSV to plain text."""
    try:
        with open(csv_file, mode='r', encoding='utf-8') as csvfile, open(txt_file, "w", encoding="utf-8") as txt:
            reader = csv.reader(csvfile)
            for row in reader:
                txt.write("\t".join(row) + "\n")
    except Exception as e:
        print(f"Error converting CSV file: {e}")

def convert_xlsx_to_txt(xlsx_file, txt_file):
    """Convert XLSX to plain text."""
    try:
        workbook = load_workbook(xlsx_file, read_only=True)
        text_content = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"Sheet: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                text_content += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            text_content += "\n"
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content)
    except Exception as e:
        print(f"Error converting XLSX file: {e}")

def convert_epub_to_txt(epub_file, txt_file):
    """Convert EPUB to plain text."""
    try:
        book = epub.read_epub(epub_file)  # Read the EPUB file
        text_content = ""
        for item in book.get_items():  # Iterate through all items in the EPUB
            if item.get_type() == epub.ITEM_DOCUMENT:  # Check if the item is a document
                soup = BeautifulSoup(item.get_content(), "html.parser")  # Parse HTML content
                text_content += soup.get_text() + "\n"  # Extract and append plain text
        with open(txt_file, "w", encoding="utf-8") as txt:  # Save to .txt file
            txt.write(text_content.strip())
    except Exception as e:
        print(f"Error converting EPUB file: {e}")

def convert_mobi_to_txt(mobi_file, txt_file):
    """Convert MOBI to plain text using kindleunpack."""
    temp_dir = mobi.extract(mobi_file)  # Extracts MOBI to a temporary directory
    epub_file = os.path.join(temp_dir, "mobi7", "content.opf")  # Locate the extracted EPUB
    if os.path.exists(epub_file):
        convert_epub_to_txt(epub_file, txt_file)
    else:
        raise ValueError("Failed to extract MOBI file.")

def convert_latex_to_txt(latex_file, txt_file):
    """Convert LaTeX to plain text using textract."""
    try:
        text_content = textract.process(latex_file, encoding="utf-8").decode("utf-8")
        with open(txt_file, "w", encoding="utf-8") as txt:
            txt.write(text_content)
    except Exception as e:
        print(f"Error converting LaTeX file: {e}")

def convert_rtf_to_txt(rtf_file, txt_file):
    """Convert RTF to plain text using striprtf."""
    with open(rtf_file, "r", encoding="utf-8") as rtf:
        rtf_content = rtf.read()
    plain_text = rtf_to_text(rtf_content)
    with open(txt_file, "w", encoding="utf-8") as txt:
        txt.write(plain_text)

def convert_odt_to_txt(odt_file, txt_file):
    """Convert ODT to plain text."""
    doc = odf_load(odt_file)
    texts = ""
    for paragraph in doc.getElementsByType(text.P):
        texts += teletype.extractText(paragraph) + "\n"
    with open(txt_file, "w", encoding="utf-8") as txt:
        txt.write(texts)

def convert_html_to_txt(html_file, txt_file):
    """Convert HTML to plain text."""
    with open(html_file, "r", encoding="utf-8") as html:
        soup = BeautifulSoup(html, "html.parser")
        text_content = soup.get_text(separator="\n")
    with open(txt_file, "w", encoding="utf-8") as txt:
        txt.write(text_content)

def convert_doc_to_txt(doc_file, txt_file):
    """Convert DOC to plain text."""
    document = Document(doc_file)
    text_content = "\n".join([paragraph.text for paragraph in document.paragraphs])
    with open(txt_file, "w", encoding="utf-8") as txt:
        txt.write(text_content)

def convert_docx_to_txt(docx_file, txt_file):
    """Convert DOCX to plain text."""
    document = Document(docx_file)
    with open(txt_file, "w", encoding="utf-8") as txt:
        for paragraph in document.paragraphs:
            txt.write(paragraph.text + "\n")

def convert_pdf_to_txt(pdf_file, txt_file):
    """Convert PDF to plain text."""
    reader = PdfReader(pdf_file)
    with open(txt_file, "w", encoding="utf-8") as txt:
        for page in reader.pages:
            txt.write(page.extract_text() + "\n")

def convert_to_txt(input_file):
    # Removed unused variable 'progress_conversion'
    start_time = time.time()
    base_name = os.path.basename(os.path.splitext(input_file)[0])
    output_file = os.path.join("Pending", f"{base_name}.txt")
    ext = os.path.splitext(input_file)[-1].lower()
    try:
        if ext == ".docx":
            convert_docx_to_txt(input_file, output_file)
        elif ext == ".pdf":
            convert_pdf_to_txt(input_file, output_file)
        elif ext == ".txt":
            with open(input_file, "r", encoding="utf-8") as f_in, \
                 open(output_file, "w", encoding="utf-8") as f_out:
                f_out.write(f_in.read())
        elif ext == ".rtf":
            convert_rtf_to_txt(input_file, output_file)
        elif ext == ".odt":
            convert_odt_to_txt(input_file, output_file)
        elif ext == ".html":
            convert_html_to_txt(input_file, output_file)
        elif ext == ".doc":
            convert_doc_to_txt(input_file, output_file)
        elif ext == ".epub":
            convert_epub_to_txt(input_file, output_file)
        elif ext == ".mobi":
            convert_mobi_to_txt(input_file, output_file)
        elif ext == ".latex":
            convert_latex_to_txt(input_file, output_file)
        elif ext == ".xlsx":
            convert_xlsx_to_txt(input_file, output_file)
        elif ext == ".csv":
            convert_csv_to_txt(input_file, output_file)
        elif ext == ".ods":
            convert_ods_to_txt(input_file, output_file)
        elif ext == ".pptx":
            convert_pptx_to_txt(input_file, output_file)
        elif ext == ".odp":
            convert_odp_to_txt(input_file, output_file)
        elif ext == ".xml":
            convert_xml_to_txt(input_file, output_file)
        elif ext in [".md", ".markdown"]:
            convert_markdown_to_txt(input_file, output_file)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        print(f"Error converting file: {e}")
    elapsed_time = time.time() - start_time
    print(f"File conversion for '{input_file}' took {elapsed_time:.2f} seconds.")

# Global variables
uploaded_files = []
plagiarism_results = []

def upload_files():
    global uploaded_files
    filepaths = askopenfilenames(filetypes=[("All files", "*.*")])
    if not filepaths:
        return

    # Show progress bar
    progress_upload = show_progress_bar(window, row=2, column=0, width=400, height=10)
    uploaded_files = []
    total_files = len(filepaths)

    try:
        for idx, filepath in enumerate(filepaths):
            base_name = os.path.basename(filepath)
            txt_path = os.path.join("Pending", f"{os.path.splitext(base_name)[0]}.txt")
            convert_to_txt(filepath)  # Convert each file to text
            uploaded_files.append(txt_path)
            progress_upload.set((idx + 1) / total_files)  # Update progress
            window.update_idletasks()  # Ensure GUI updates

        messagebox.showinfo("Upload Complete", "Files uploaded and converted successfully.")
    finally:
        # Hide progress bar after completion
        hide_progress_bar(progress_upload)

def extract_copied_texts(file1, file2, min_similarity=0.7):
    """Extract copied texts between two files."""
    with open(file1, encoding='utf-8') as f1, open(file2, encoding='utf-8') as f2:
        text1 = f1.read()
        text2 = f2.read()
    matcher = SequenceMatcher(None, text1, text2)
    copied_texts = []
    for match in matcher.get_matching_blocks():
        if match.size > 10:
            copied_text = text1[match.a:match.a + match.size].strip()
            if copied_text:
                block_similarity = matcher.real_quick_ratio()
                if block_similarity >= min_similarity:
                    copied_texts.append({
                        "source_document_1": os.path.basename(file1),
                        "source_document_2": os.path.basename(file2),
                        "copied_text": copied_text,
                        "similarity": block_similarity
                    })
    if not copied_texts:
        embeds = MODEL.encode([text1, text2])
        overall_similarity = cosine_similarity([embeds[0]], [embeds[1]])[0][0]
        if overall_similarity >= min_similarity:
            copied_texts.append({
                "source_document_1": os.path.basename(file1),
                "source_document_2": os.path.basename(file2),
                "copied_text": "(No specific copied text detected, but overall similarity is high.)",
                "similarity": overall_similarity
            })
    return copied_texts

def check_uploaded_files_plagiarism():
    global uploaded_files, plagiarism_results
    if not uploaded_files:
        messagebox.showwarning("No Files", "No files have been uploaded.")
        return

    # Show progress bar
    progress_check = show_progress_bar(window, row=4, column=0, width=400, height=10)
    progress_check.set(0)

    try:
        corrupt_folder = os.path.join("Pending", "corrupt documents")
        os.makedirs(corrupt_folder, exist_ok=True)
        valid_files = []
        corrupted_files = []
        file_contents = []

        # Step 1: Load and validate files
        total_steps = 4  # Embedding, clustering, comparison, logging
        step = 1
        for file in uploaded_files:
            try:
                with open(file, encoding='utf-8') as f:
                    content = f.read().strip()
                if not content:
                    file_name = os.path.basename(file)
                    corrupt_file_path = os.path.join(corrupt_folder, file_name)

                    # Handle file conflicts by appending a timestamp
                    while os.path.exists(corrupt_file_path):
                        base_name, ext = os.path.splitext(file_name)
                        timestamp = time.strftime("%Y%m%d%H%M%S")
                        file_name = f"{base_name}_{timestamp}{ext}"
                        corrupt_file_path = os.path.join(corrupt_folder, file_name)

                    os.rename(file, corrupt_file_path)
                    corrupted_files.append(file_name)
                    continue
                file_contents.append(content)
                valid_files.append(file)
            except Exception as e:
                file_name = os.path.basename(file)
                corrupt_file_path = os.path.join(corrupt_folder, file_name)

                # Handle file conflicts by appending a timestamp
                while os.path.exists(corrupt_file_path):
                    base_name, ext = os.path.splitext(file_name)
                    timestamp = time.strftime("%Y%m%d%H%M%S")
                    file_name = f"{base_name}_{timestamp}{ext}"
                    corrupt_file_path = os.path.join(corrupt_folder, file_name)

                os.rename(file, corrupt_file_path)
                corrupted_files.append(file_name)
        if not file_contents:
            messagebox.showerror("Empty Files", "All uploaded files are empty, corrupted, or contain no meaningful content.")
            return
        progress_check.set(step / total_steps)  # Update progress
        window.update_idletasks()

        # Step 2: Generate embeddings
        embeddings_start = time.time()
        embeddings = MODEL.encode(file_contents)
        embedding_time = time.time() - embeddings_start
        step += 1
        progress_check.set(step / total_steps)  # Update progress
        window.update_idletasks()

        # Step 3: Cluster optimization
        cluster_start = time.time()
        optimal_clusters = optimal_cluster_count(embeddings)
        kmeans = MiniBatchKMeans(n_clusters=optimal_clusters, random_state=42, batch_size=100)
        clusters = kmeans.fit_predict(embeddings)
        clustering_time = time.time() - cluster_start
        step += 1
        progress_check.set(step / total_steps)  # Update progress
        window.update_idletasks()

        # Step 4: Group documents by cluster and compare
        cluster_groups = defaultdict(list)
        for idx, cluster_id in enumerate(clusters):
            cluster_groups[cluster_id].append((valid_files[idx], embeddings[idx]))

        compare_start = time.time()
        results = Parallel(n_jobs=-1, prefer="threads")(
            delayed(cluster_comparison)(cluster) for cluster in cluster_groups.values()
        )
        comparison_time = time.time() - compare_start
        plagiarism_results = [item for sublist in results for item in sublist]

        step += 1
        progress_check.set(step / total_steps)  # Update progress
        window.update_idletasks()

        total_time = time.time() - embeddings_start
        print(f"Plagiarism check completed. Embedding time: {embedding_time:.2f}s, Clustering time: {clustering_time:.2f}s, Comparison time: {comparison_time:.2f}s, Total time: {total_time:.2f}s")

        if corrupted_files:
            messagebox.showinfo(
                "Plagiarism Check Complete",
                f"Plagiarism check completed successfully.\nThe following files were corrupted or empty and moved to 'corrupt documents':\n{', '.join(corrupted_files)}"
            )
        else:
            messagebox.showinfo("Plagiarism Check Complete", "Plagiarism check completed successfully.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred during the plagiarism check: {e}")
    finally:
        # Hide progress bar after completion
        hide_progress_bar(progress_check)

def optimal_cluster_count(embeddings, max_clusters=5):
    """Determine the optimal number of clusters."""
    best_score = -1
    optimal = 2
    for n in range(2, min(max_clusters, len(embeddings)) + 1):
        kmeans = MiniBatchKMeans(n_clusters=n, random_state=42, batch_size=100)
        kmeans.fit(embeddings)
        score = silhouette_score(embeddings, kmeans.labels_)
        if score > best_score:
            best_score = score
            optimal = n
    return optimal

def cluster_comparison(cluster_docs):
    """Compare documents within a cluster."""
    filenames = [doc[0] for doc in cluster_docs]
    embeds = np.array([doc[1] for doc in cluster_docs])
    sim_matrix = cosine_similarity(embeds)
    results = []
    for i in range(sim_matrix.shape[0]):
        for j in range(i + 1, sim_matrix.shape[1]):
            if sim_matrix[i, j] >= 0.7:
                pair = sorted([os.path.basename(filenames[i]), os.path.basename(filenames[j])])
                # Convert similarity score to percentage
                similarity_percentage = int(sim_matrix[i, j] * 100)  # Convert to integer percentage
                results.append({
                    "Assignment 1": pair[0],
                    "Assignment 2": pair[1],
                    "Similarity Score": f"{similarity_percentage}%",  # Format as percentage
                    "Plagiarism Status": "Flagged"
                })
    return results

def extract_keywords(text, num_keywords=5):
    """Extract keywords from text."""
    words = word_tokenize(text.lower())
    stop_words = set(stopwords.words('english'))
    filtered_words = [word for word in words if word.isalpha() and word not in stop_words]
    freq_dist = nltk.FreqDist(filtered_words)
    return [word for word, _ in freq_dist.most_common(num_keywords)]

def show_copied_texts():
    if not plagiarism_results:
        messagebox.showerror("No Results", "No plagiarism results to display.")
        return

    result_window = customtkinter.CTkToplevel()
    result_window.title("Plagiarism Results Window")
    result_window.geometry("1200x800")
    result_window.state("zoomed")

    main_frame = customtkinter.CTkFrame(result_window)
    main_frame.pack(fill="both", expand=True, padx=10, pady=10)

    text_widget = tk.Text(main_frame, wrap="word", font=("Consolas", 12))
    scrollbar = customtkinter.CTkScrollbar(main_frame, command=text_widget.yview)
    text_widget.configure(yscrollcommand=scrollbar.set)
    text_widget.grid(row=0, column=0, sticky="nsew")
    scrollbar.grid(row=0, column=1, sticky="ns")

    main_frame.grid_rowconfigure(0, weight=1)
    main_frame.grid_columnconfigure(0, weight=1)

    # Header for the results
    text_widget.insert("end", "\nPlagiarism Results:\n")
    text_widget.insert("end", "{:<10} {:<30} {:<30} {:<20}\n".format(
        "ID", "Source Document", "Copied Document", "Similarity Score"
    ))
    text_widget.insert("end", "-" * 100 + "\n")
    id_counter = 1

    for result in plagiarism_results:
        if result["Plagiarism Status"] == "Flagged":
            source_doc = result["Assignment 1"]
            copied_doc = result["Assignment 2"]
            sim_score = result["Similarity Score"]  # Already formatted as percentage
            text_widget.insert("end", "{:<10} {:<30} {:<30} {}\n".format(
                id_counter, source_doc, copied_doc, sim_score
            ))

            file1 = os.path.join("Pending", source_doc)
            file2 = os.path.join("Pending", copied_doc)
            copied_texts = extract_copied_texts(file1, file2)
            if copied_texts:
                text_widget.insert("end", "\nHighlighted Copied Text:\n")
                for copied_text in copied_texts:
                    text_widget.insert("end", f"{copied_text['copied_text']}\n")
            else:
                with open(file1, encoding='utf-8') as f1, open(file2, encoding='utf-8') as f2:
                    text1 = f1.read()
                    text2 = f2.read()
                combined_text = text1 + " " + text2
                keywords = extract_keywords(combined_text, num_keywords=5)
                text_widget.insert("end", "\nExtracted Keywords:\n")
                if keywords:
                    text_widget.insert("end", ", ".join(keywords) + "\n")
                else:
                    text_widget.insert("end", "(No keywords could be extracted.)\n")

            text_widget.insert("end", "-" * 100 + "\n")
            id_counter += 1

    # Dynamically configure text widget appearance based on appearance mode
    current_mode = customtkinter.get_appearance_mode()
    if current_mode == "Dark":
        text_widget.configure(fg="white", bg="black")  # Dark mode colors
    else:
        text_widget.configure(fg="black", bg="white")  # Light mode colors

    text_widget.configure(state="disabled")  # Make the text widget read-only

    # Close button
    close_btn = customtkinter.CTkButton(
        main_frame,
        text="Close Window",
        command=result_window.destroy,
        fg_color="#4CAF50",
        hover_color="#45a049"
    )
    close_btn.grid(row=1, column=0, pady=10)

def export_report():
    """Export plagiarism results to CSV and PDF."""
    if not plagiarism_results:
        messagebox.showerror("No Results", "No plagiarism results to export.")
        return
    flagged_results = [result for result in plagiarism_results if result["Plagiarism Status"] == "Flagged"]
    if not flagged_results:
        messagebox.showinfo("No Flagged Documents", "No flagged documents to export.")
        return
    simplified_results = [
        {
            "Source Document 1": result["Assignment 1"],
            "Source Document 2": result["Assignment 2"],
            "Similarity Score": result["Similarity Score"],
            "Plagiarism Status": result["Plagiarism Status"]
        }
        for result in flagged_results
    ]
    csv_filename = os.path.join("Reports", "plagiarism_report.csv")
    keys = simplified_results[0].keys()
    with open(csv_filename, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=keys)
        writer.writeheader()
        writer.writerows(simplified_results)
    pdf_filename = os.path.join("Reports", "plagiarism_report.pdf")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, txt="Plagiarism Report", ln=1, align='C')
    pdf.ln(10)
    for result in simplified_results:
        pdf.multi_cell(0, 10, txt=f"Source Document 1: {result['Source Document 1']}\n"
                                 f"Source Document 2: {result['Source Document 2']}\n"
                                 f"Similarity Score: {result['Similarity Score']}\n"
                                 f"Plagiarism Status: {result['Plagiarism Status']}\n")
        pdf.ln(5)
    pdf.output(pdf_filename)
    messagebox.showinfo("Report Generated", f"Report saved successfully in Reports folder.\nCSV: {csv_filename}\nPDF: {pdf_filename}")

def open_report():
    """Open the generated report."""
    # Create a custom dialog box with "PDF" and "CSV" buttons
    choice = messagebox.askyesnocancel("Open Report", "Would you like to open the report in PDF format? Select 'PDF' or 'CSV'.")
    
    if choice is None:  # User clicked Cancel
        return
    
    if choice:  # User selected PDF
        pdf_path = os.path.join("Reports", "plagiarism_report.pdf")
        if os.path.exists(pdf_path):
            os.startfile(pdf_path)  # Open PDF
        else:
            messagebox.showerror("File Not Found", "The PDF report does not exist.")
    else:  # User selected CSV
        csv_path = os.path.join("Reports", "plagiarism_report.csv")
        if os.path.exists(csv_path):
            os.startfile(csv_path)  # Open CSV
        else:
            messagebox.showerror("File Not Found", "The CSV report does not exist.")

def reset_application():
    """Reset the application state after user confirmation."""
    confirm_reset = messagebox.askyesno("Confirm Reset", "Are you sure you want to reset the application? This will clear all uploaded files and results.")
    if confirm_reset:
        global uploaded_files, plagiarism_results
        uploaded_files = []
        plagiarism_results = []

        # Remove files from "Pending" folder
        for file in os.listdir("Pending"):
            file_path = os.path.join("Pending", file)
            if os.path.isfile(file_path):
                os.remove(file_path)

        # Confirmation message after reset
        messagebox.showinfo("Reset Complete", "Application data has been reset successfully.")

def toggle_appearance_mode():
    """Toggle light/dark mode."""
    global btn_toggle_mode
    current_mode = customtkinter.get_appearance_mode()
    new_mode = "Light" if current_mode == "Dark" else "Dark"
    customtkinter.set_appearance_mode(new_mode)

    # Update the tooltip text
    tooltip_text = "Switch to Light Mode" if new_mode == "Dark" else "Switch to Dark Mode"
    Tooltip(btn_toggle_mode, tooltip_text).hide_tooltip()  # Reapply tooltip

    # Update widget colors
    update_widget_colors()

def update_widget_colors():
    """Update widget colors based on appearance mode."""
    current_mode = customtkinter.get_appearance_mode()
    if current_mode == "Light":
        window.configure(fg_color="#F5F5F5")
        welcome_frm.configure(fg_color="#F5F5F5")
        welcome_lbl.configure(text_color="#333333")
        button_frm.configure(fg_color="#F5F5F5")
        for button in [btn_upload, btn_check, btn_report, btn_open_report, btn_reset, btn_show_copied_texts]:
            button.configure(fg_color="#4CAF50", hover_color="#45a049", text_color="#FFFFFF")
        btn_reset.configure(fg_color="#FF0000", hover_color="#CC0000")
        btn_toggle_mode.configure(fg_color="#F5F5F5", hover_color="#F5F5F5", text_color="#000000")
    else:
        window.configure(fg_color="#2E2E2E")
        welcome_frm.configure(fg_color="#2E2E2E")
        welcome_lbl.configure(text_color="yellow")
        button_frm.configure(fg_color="#2E2E2E")
        for button in [btn_upload, btn_check, btn_report, btn_open_report, btn_reset, btn_show_copied_texts]:
            button.configure(fg_color="#4CAF50", hover_color="#45a049", text_color="#FFFFFF")
        btn_reset.configure(fg_color="#FF0000", hover_color="#CC0000")
        btn_toggle_mode.configure(fg_color="#2E2E2E", hover_color="#2E2E2E", text_color="#FFFFFF")

class Tooltip:
    """Create a modern, styled tooltip for a given widget."""
    def __init__(self, widget, text):
        self.widget = widget
        self.text = text
        self.tooltip_window = None
        self.widget.bind("<Enter>", self.show_tooltip)
        self.widget.bind("<Leave>", self.hide_tooltip)

    def show_tooltip(self, event=None):
        """Display the tooltip."""
        if self.tooltip_window or not self.text:
            return

        # Determine position
        x, y, _, _ = self.widget.bbox("insert")
        if x is None or y is None:
            x, y = self.widget.winfo_pointerxy()
        x += self.widget.winfo_rootx() + 25
        y += self.widget.winfo_rooty() + 25

        # Create the tooltip window
        self.tooltip_window = tk.Toplevel(self.widget)
        self.tooltip_window.wm_overrideredirect(True)  # Remove window decorations
        self.tooltip_window.wm_geometry(f"+{x}+{y}")

        # Get current appearance mode
        current_mode = customtkinter.get_appearance_mode()

        # Set background and foreground colors based on mode
        bg_color = "#F5F5F5" if current_mode == "Light" else "#2E2E2E"
        fg_color = "#333333" if current_mode == "Light" else "#FFFFFF"

        # Add a label to the tooltip window
        label = customtkinter.CTkLabel(
            self.tooltip_window,
            text=self.text,
            bg_color=bg_color,
            text_color=fg_color,
            font=("Arial", 10),
            corner_radius=8,
            padx=10,
            pady=5
        )
        label.pack()

    def hide_tooltip(self, event=None):
        """Hide the tooltip."""
        if self.tooltip_window:
            self.tooltip_window.destroy()
            self.tooltip_window = None

# GUI implementation
window = customtkinter.CTk()
customtkinter.set_appearance_mode("dark")
window.title("Internal Plagiarism Detector")
window.state("zoomed")

# Load icons
if os.path.exists("sun.png"):
    sun_icon = customtkinter.CTkImage(light_image=Image.open("sun.png"), size=(32, 32))
else:
    sun_icon = None

if os.path.exists("moon.png"):
    moon_icon = customtkinter.CTkImage(dark_image=Image.open("moon.png"), size=(32, 32))
else:
    moon_icon = None

# GUI components
welcome_frm = customtkinter.CTkFrame(window)
welcome_msg_variable = tk.StringVar(welcome_frm, "Welcome to the state-of-the-art Internal Plagiarism Detector")
welcome_lbl = customtkinter.CTkLabel(welcome_frm, textvariable=welcome_msg_variable,
                                     height=100, corner_radius=20, 
                                     text_color="yellow", font=("Comic Sans MS bold", 30))
welcome_lbl.grid(row=0, column=0, padx=200, pady=(20, 0), sticky="nsew")

button_frm = customtkinter.CTkFrame(window)

# UPLOAD FILES Button
btn_upload = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Green", text="UPLOAD FILES", command=upload_files)
btn_upload.grid(row=0, column=0, padx=(180, 10), pady=10)
Tooltip(btn_upload, "Upload files to check for plagiarism.")

# CHECK PLAGIARISM Button
btn_check = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Green", text="CHECK PLAGIARISM", command=check_uploaded_files_plagiarism)
btn_check.grid(row=0, column=1, padx=10, pady=10)
Tooltip(btn_check, "Check uploaded files for plagiarism.")

# EXPORT REPORT Button
btn_report = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Green", text="EXPORT REPORT", command=export_report)
btn_report.grid(row=0, column=2, padx=10, pady=10)
Tooltip(btn_report, "Export plagiarism results to CSV and PDF formats.")

# OPEN REPORT Button
btn_open_report = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Green", text="OPEN REPORT", command=open_report)
btn_open_report.grid(row=0, column=3, padx=10, pady=10)
Tooltip(btn_open_report, "Open the generated plagiarism report in PDF or CSV format.")

# RESET Button
btn_reset = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Red", text="RESET", command=reset_application)
btn_reset.grid(row=0, column=4, padx=10, pady=10)
Tooltip(btn_reset, "Reset the application and clear all uploaded files and results.")

# SHOW COPIED TEXTS Button
btn_show_copied_texts = customtkinter.CTkButton(button_frm, corner_radius=30, hover_color="Green", text="SHOW COPIED TEXTS", command=show_copied_texts)
btn_show_copied_texts.grid(row=0, column=5, padx=10, pady=10)
Tooltip(btn_show_copied_texts, "Show detailed copied texts between flagged documents.")

button_frm.grid(row=1, column=0, padx=20, pady=(10, 50), sticky="nsew")

# Toggle Appearance Mode Button
top_right_frame = customtkinter.CTkFrame(window)
top_right_frame.grid(row=0, column=1, padx=20, pady=20, sticky="ne")
btn_toggle_mode = customtkinter.CTkButton(top_right_frame, text="", image=sun_icon if sun_icon else None, 
                                          command=toggle_appearance_mode, fg_color="#2E2E2E", hover_color="#2E2E2E", 
                                          width=40, height=40)
btn_toggle_mode.grid(row=0, column=0)
Tooltip(btn_toggle_mode, "Toggle light/dark mode.")

welcome_frm.grid(row=0, column=0, padx=20, pady=(50, 10), sticky="nsew")
button_frm.grid(row=1, column=0, padx=20, pady=(10, 50), sticky="nsew")
window.grid_rowconfigure(0, weight=0)
window.grid_columnconfigure(0, weight=1)

update_widget_colors()
window.mainloop()